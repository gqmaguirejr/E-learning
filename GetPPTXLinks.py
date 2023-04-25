#!/usr/bin/env python3
# -*- coding: utf-8 -*-
# -*- mode: python; python-indent-offset: 4 -*-
#
# ./GetPPTXLinks.py path [options]
#
# Purpose: The program takes in a PPTX file and extracts the links
#
# Output: outputs a spreadsheet
#
# Example:
# ./GetPPTXLinks.py Lecture-2-2-dl-backprop1.pptx -v
#
# ./GetPPTXLinks.py /z3/maguire/Nvidia/DeepLearningKit/ -r -o DeepLearningKit_PPTX_Links.xlsx
#
#
# uses wildcard parsing from https://github.com/Colin-Fredericks/hx_util/blob/master/hx_util/GetWordLinks.py
# that was last updated on April 3rd 2019
# 
# 2023-04-25 G. Q. Maguire Jr.
#
import re
import sys
import subprocess

import json
import argparse
import os			# to make OS calls, here to get time zone info

import time

import pprint

from collections import defaultdict


import datetime
import isodate                  # for parsing ISO 8601 dates and times
import pytz                     # for time zones
from dateutil.tz import tzlocal

# for dealing with the PPTX file - which is a ZIP file
import zipfile

# for creating files and directories
from pathlib import Path

# for UNIX styte path name expansion (i.e., wildcards)
import glob

# Colin-Fredericks used BeautifulSoup and unicodecsv
# from bs4 import BeautifulSoup
# import unicodecsv as csv  # https://pypi.python.org/pypi/unicodecsv/0.14.1

# instead we will use python-pptx
from pptx import Presentation

# use Pandas to make the spreadsheet
import pandas as pd

try:
    import zlib
    compression = zipfile.ZIP_DEFLATED
except:
    compression = zipfile.ZIP_STORED

modes = { zipfile.ZIP_DEFLATED: 'deflated',
          zipfile.ZIP_STORED:   'stored',
          }

instructions = """
Usage:

python3 GetWordLinks.py path/to/file/ (options)

Extract all hyperlinks from a .docx or .docm file,
including link destination and linked text,
and store them in a .xlsx file.
If you feed it a folder, it includes all the files in the folder.

Options:
  -h  Print this message and quit.
  -r  Recursive - includes nested folders.
  -o  Set an output filename as the next argument.
"""

def getLinks(filename, args, dirpath):
    global Verbose_Flag

    links_with_urls=[]

    # Open the .pptx file as if it were a zip (because it is)
    fullname = os.path.join(dirpath or "", filename)
    if Verbose_Flag:
        print("getting links for file: {}".format(fullname))
    try:
        prs = Presentation(fullname)
    except:
        print("Error encountered with processing: {}".format(fullname))
        return links_with_urls
        pass
    
    # text_runs will be populated with a list of strings,
    # one for each text run in presentation
    #for slide in prs.slides:
    for slideIdx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    try:
                        link=run.hyperlink.address
                    except:
                        pass
                    if link:
                        links_with_urls.append({'link': link, 'text': run.text, 'filename': os.path.basename(fullname), 'slide': slideIdx})

    # Return a list of dicts full of link info
    return links_with_urls


def main(args):
    global Verbose_Flag
    
    print("Getting .pptx Links")

    # Handle arguments and flags
    argp = argparse.ArgumentParser(usage=instructions, add_help=False)
    argp.add_argument("--help", "-h", action="store_true")
    argp.add_argument("-r", action="store_true")
    argp.add_argument("-o", action="store")
    argp.add_argument("file_names", nargs="*")
    argp.add_argument('-v', '--verbose', required=False,
                      default=False,
                      action="store_true",
                      help="Print lots of output to stdout")

    args = argp.parse_args(args)
    Verbose_Flag=vars(args)["verbose"]
    # Replace arguments with wildcards with their expansion.
    # If a string does not contain a wildcard, glob will return it as is.
    # Mostly important if we run this on Windows systems.
    file_names = list()

    for name in args.file_names:
        file_names += glob.glob(glob.escape(name))

    # If the filenames don't exist, say so and quit.
    if file_names == []:
        sys.exit("No file or directory found by that name.")

    # Don't run the script on itself.
    if sys.argv[0] in file_names:
        file_names.remove(sys.argv[0])

    if args.help:
        sys.exit(instructions)

    filecount = 0
    linklist = []
    target_is_folder = False

    for name in file_names:
        # Make sure single files exist.
        assert os.path.exists(name), "File or directory not found."

        # If it's just a file...
        if os.path.isfile(name):
            # Make sure this is a Word file (just check extension)
            if name.lower().endswith(".pptx"):
                # Get links from that file.
                linklist.extend(getLinks(name, args, False))
                filecount += 1

        # If it's a directory:
        if os.path.isdir(name):
            target_is_folder = True
            # Recursive version using os.walk for all levels.
            if args.r:
                for dirpath, dirnames, files in os.walk(name):
                    for eachfile in files:
                        # Get links for every file in that directory.
                        if eachfile.lower().endswith(".pptx"):
                            linklist.extend(getLinks(eachfile, args, dirpath))
                            filecount += 1
            # Non-recursive version breaks os.walk after the first level.
            else:
                topfiles = []
                for (dirpath, dirnames, files) in os.walk(name):
                    topfiles.extend(files)
                    break
                for eachfile in topfiles:
                    if eachfile.lower().endswith(".pptx"):
                        linklist.extend(getLinks(eachfile, args, dirpath))
                        filecount += 1

    # Otherwise, output a file and print some info.
    print(
        "\nChecked "
        + str(filecount)
        + " .pptx file"
        + ("s" if filecount > 1 else "")
        + " for links."
    )

    # Create output file as sibling to the original target of the script.
    outFileName = args.o if args.o else "PPTX_Links.xlsx"
    if target_is_folder:
        outFileFolder = os.path.abspath(os.path.join(file_names[0], os.pardir))
        outFilePath = os.path.join(outFileFolder, outFileName)
    else:
        outFilePath = os.path.join(os.path.dirname(file_names[0]), outFileName)

    if Verbose_Flag:
        print(linklist)

    if (linklist):
        linklist_df=pd.json_normalize(linklist)

        if Verbose_Flag:
            print(f'{outFileName=}')
            print(f'{outFilePath=}')
        writer = pd.ExcelWriter(outFilePath, engine='xlsxwriter')
        linklist_df.to_excel(writer, sheet_name='Links')

        # Close the Pandas Excel writer and output the Excel file.
        writer.close()

    print("Spreadsheet created: " + outFileName)
    print("Location: " + outFilePath)


if __name__ == '__main__':
    # this won't be run when imported
    sys.exit(main(sys.argv[1:]))
